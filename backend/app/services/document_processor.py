import re
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from langdetect import detect

from app.config import CHUNK_SIZE, CHUNK_OVERLAP


def extract_pdf(path: str) -> str:
    doc = fitz.open(path)
    pages = []
    for page in doc:
        text = page.get_text()
        text = _clean_page_text(text)
        if text.strip():
            pages.append(text)
    doc.close()
    return "\n\n".join(pages)


def extract_docx(path: str) -> str:
    doc = DocxDocument(path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def extract_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _clean_page_text(text: str) -> str:
    text = re.sub(r"^[\d\s]*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"(Page\s*\d+|^\d+\s*$)", "", text, flags=re.MULTILINE | re.IGNORECASE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    if len(words) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start = end - overlap

    return chunks


def detect_language(text: str) -> str:
    try:
        sample = text[:3000]
        lang = detect(sample)
        return lang
    except Exception:
        return "en"


EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
}


def process_document(file_path: str, file_type: str) -> tuple[str, list[str], str]:
    extractor = EXTRACTORS.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")

    raw_text = extractor(file_path)
    if not raw_text.strip():
        raise ValueError("No text could be extracted from the document")

    language = detect_language(raw_text)
    chunks = chunk_text(raw_text)

    return raw_text, chunks, language
